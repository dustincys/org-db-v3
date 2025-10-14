"""Service for converting linked files to markdown using lightweight libraries."""
import logging
import hashlib
from pathlib import Path
from typing import Optional, Dict, Any

logger = logging.getLogger(__name__)


class DoclingService:
    """Service for converting documents to markdown using lightweight libraries.

    Supports:
    - PDF files via pymupdf4llm
    - DOCX files via python-docx
    - PPTX files via python-pptx

    Other formats are not supported and will be logged.
    """

    # Supported file extensions (only lightweight libraries)
    SUPPORTED_EXTENSIONS = {
        '.pdf',   # pymupdf4llm
        '.docx',  # python-docx
        '.pptx',  # python-pptx
    }

    def __init__(self):
        """Initialize the document conversion service."""
        logger.info("DoclingService initialized with lightweight libraries (PDF, DOCX, PPTX)")

    @staticmethod
    def calculate_md5(file_path: str) -> str:
        """Calculate MD5 hash of a file for change detection."""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def is_supported(self, file_path: str) -> bool:
        """Check if a file extension is supported (PDF, DOCX, PPTX only)."""
        ext = Path(file_path).suffix.lower()
        is_supported = ext in self.SUPPORTED_EXTENSIONS
        if not is_supported and ext:
            logger.warning(f"File type '{ext}' is not supported yet (only PDF, DOCX, PPTX)")
        return is_supported

    def convert_to_markdown(
        self,
        file_path: str,
        max_file_size: int = 52428800  # 50MB default
    ) -> Dict[str, Any]:
        """
        Convert a document to markdown using docling.

        Args:
            file_path: Path to the file to convert
            max_file_size: Maximum file size in bytes (default 50MB)

        Returns:
            Dictionary with:
                - status: 'success', 'error', 'skipped'
                - markdown: The converted markdown text (if successful)
                - md5: MD5 hash of the source file
                - error: Error message (if failed)
                - file_size: Size of the file in bytes
        """
        path = Path(file_path)

        # Check file exists
        if not path.exists():
            logger.warning(f"File not found: {file_path}")
            return {
                "status": "error",
                "error": "File not found",
                "md5": None,
                "file_size": 0
            }

        # Check file size
        file_size = path.stat().st_size
        if file_size > max_file_size:
            logger.warning(f"File too large: {file_path} ({file_size} bytes)")
            return {
                "status": "skipped",
                "error": f"File too large: {file_size / 1024 / 1024:.1f}MB",
                "md5": None,
                "file_size": file_size
            }

        # Check extension
        if not self.is_supported(file_path):
            ext = path.suffix
            logger.warning(f"Unsupported file extension: {ext}")
            return {
                "status": "error",
                "error": f"Unsupported file extension: {ext}",
                "md5": None,
                "file_size": file_size
            }

        # Calculate MD5
        try:
            md5 = self.calculate_md5(file_path)
        except Exception as e:
            logger.error(f"Error calculating MD5 for {file_path}: {e}")
            return {
                "status": "error",
                "error": f"Error calculating MD5: {str(e)}",
                "md5": None,
                "file_size": file_size
            }

        # Convert document based on extension
        ext = path.suffix.lower()

        if ext == '.pdf':
            return self._convert_pdf_with_pymupdf(file_path, md5, file_size)
        elif ext == '.docx':
            return self._convert_docx_with_python_docx(file_path, md5, file_size)
        elif ext == '.pptx':
            return self._convert_pptx_with_python_pptx(file_path, md5, file_size)
        else:
            # Unsupported format
            logger.warning(f"Cannot convert {file_path}: unsupported file type '{ext}'")
            return {
                "status": "error",
                "error": f"Unsupported file type: {ext} (only PDF, DOCX, PPTX supported)",
                "md5": md5,
                "file_size": file_size
            }

    def _convert_pdf_with_pymupdf(
        self,
        file_path: str,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert PDF using lightweight pymupdf4llm (fast, low memory)."""
        try:
            import pymupdf4llm

            logger.info(f"Converting PDF {file_path} with pymupdf4llm...")
            markdown_text = pymupdf4llm.to_markdown(file_path)

            logger.info(f"Successfully converted {file_path} ({len(markdown_text)} chars)")
            return {
                "status": "success",
                "markdown": markdown_text,
                "md5": md5,
                "file_size": file_size
            }

        except Exception as e:
            logger.error(f"Error converting PDF {file_path}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }

    def _convert_docx_with_python_docx(
        self,
        file_path: str,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert DOCX using lightweight python-docx (fast, low memory)."""
        try:
            from docx import Document

            logger.info(f"Converting DOCX {file_path} with python-docx...")
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])

            logger.info(f"Successfully converted {file_path} ({len(text)} chars)")
            return {
                "status": "success",
                "markdown": text,  # Plain text, not markdown, but good enough
                "md5": md5,
                "file_size": file_size
            }

        except Exception as e:
            logger.error(f"Error converting DOCX {file_path}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }

    def _convert_pptx_with_python_pptx(
        self,
        file_path: str,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert PPTX using lightweight python-pptx (fast, low memory)."""
        try:
            from pptx import Presentation

            logger.info(f"Converting PPTX {file_path} with python-pptx...")
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

            result_text = "\n".join(text)

            logger.info(f"Successfully converted {file_path} ({len(result_text)} chars)")
            return {
                "status": "success",
                "markdown": result_text,  # Plain text, not markdown, but good enough
                "md5": md5,
                "file_size": file_size
            }

        except Exception as e:
            logger.error(f"Error converting PPTX {file_path}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }



# Global instance
_docling_service: Optional[DoclingService] = None


def get_docling_service() -> DoclingService:
    """Get or create the global document conversion service instance."""
    global _docling_service
    if _docling_service is None:
        _docling_service = DoclingService()
    return _docling_service
